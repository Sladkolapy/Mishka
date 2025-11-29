from fastapi import FastAPI, APIRouter, HTTPException, Depends, UploadFile, File, Form, Query
from fastapi.responses import FileResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from pydantic import BaseModel, Field, EmailStr
from typing import List, Optional
import uuid
from datetime import datetime, timezone, timedelta
import bcrypt
import jwt
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
import shutil
import json
import re
import copy
from urllib.parse import quote

# Document processing libraries
import openpyxl
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
import PyPDF2
import xlrd
import io

# AI Integration
from emergentintegrations.llm.chat import LlmChat, UserMessage

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ.get('DB_NAME', 'docai_chat')]

# JWT Configuration
JWT_SECRET = os.environ.get('JWT_SECRET', 'your-super-secret-jwt-key-change-in-production')
JWT_ALGORITHM = 'HS256'
JWT_EXPIRATION_HOURS = 24 * 7

# Emergent LLM Key
EMERGENT_LLM_KEY = os.environ.get('EMERGENT_LLM_KEY', '')

# Admin email
ADMIN_EMAIL = 'mishkasladkolapka@gmail.com'

# File storage
UPLOAD_DIR = ROOT_DIR / 'uploads'
UPLOAD_DIR.mkdir(exist_ok=True)
GENERATED_DIR = ROOT_DIR / 'generated'
GENERATED_DIR.mkdir(exist_ok=True)

# Token costs
TOKEN_COSTS = {
    'pptx_create': 65,
    'pptx_edit': 10,
    'docx_create': 35,
    'docx_edit': 6,
    'xlsx_create': 40,
    'xlsx_edit': 7,
    'analyze': 5,
}

# Initial tokens for new users
INITIAL_TOKENS = 300

app = FastAPI()
api_router = APIRouter(prefix="/api")
security = HTTPBearer(auto_error=False)

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# ==================== MODELS ====================

class UserCreate(BaseModel):
    email: EmailStr
    password: str
    agree_terms: bool = True

class UserLogin(BaseModel):
    email: EmailStr
    password: str

class UserResponse(BaseModel):
    id: str
    email: str
    balance: int
    is_admin: bool
    created_at: datetime

class TokenResponse(BaseModel):
    access_token: str
    token_type: str = "bearer"
    user: UserResponse

class ChatCreate(BaseModel):
    title: Optional[str] = None

class ChatResponse(BaseModel):
    id: str
    user_id: str
    title: str
    created_at: datetime
    updated_at: datetime

class MessageCreate(BaseModel):
    content: str

class MessageResponse(BaseModel):
    id: str
    chat_id: str
    role: str
    content: str
    file_id: Optional[str] = None
    file_name: Optional[str] = None
    tokens_used: Optional[int] = None
    created_at: datetime

class FileResponse_(BaseModel):
    id: str
    filename: str
    file_type: str
    is_generated: bool
    created_at: datetime

class ChatDetailResponse(BaseModel):
    id: str
    user_id: str
    title: str
    messages: List[MessageResponse]
    files: List[FileResponse_]
    created_at: datetime
    updated_at: datetime

class TopUpRequest(BaseModel):
    amount: int

class BalanceResponse(BaseModel):
    balance: int
    token_costs: dict

class AdminUserUpdate(BaseModel):
    balance: Optional[int] = None
    is_blocked: Optional[bool] = None

# ==================== HELPERS ====================

def hash_password(password: str) -> str:
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def verify_password(password: str, hashed: str) -> bool:
    return bcrypt.checkpw(password.encode('utf-8'), hashed.encode('utf-8'))

def create_token(user_id: str) -> str:
    payload = {
        'user_id': user_id,
        'exp': datetime.now(timezone.utc) + timedelta(hours=JWT_EXPIRATION_HOURS)
    }
    return jwt.encode(payload, JWT_SECRET, algorithm=JWT_ALGORITHM)

async def get_current_user(credentials: HTTPAuthorizationCredentials = Depends(security)):
    if not credentials:
        raise HTTPException(status_code=401, detail='Требуется авторизация')
    try:
        payload = jwt.decode(credentials.credentials, JWT_SECRET, algorithms=[JWT_ALGORITHM])
        user_id = payload.get('user_id')
        user = await db.users.find_one({'id': user_id})
        if not user:
            raise HTTPException(status_code=401, detail='Пользователь не найден')
        if user.get('is_blocked'):
            raise HTTPException(status_code=403, detail='Аккаунт заблокирован')
        return user
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail='Токен истёк')
    except jwt.InvalidTokenError:
        raise HTTPException(status_code=401, detail='Неверный токен')

async def get_admin_user(user = Depends(get_current_user)):
    if not user.get('is_admin'):
        raise HTTPException(status_code=403, detail='Требуются права администратора')
    return user

def is_admin(email: str) -> bool:
    return email.lower() == ADMIN_EMAIL.lower()

async def deduct_tokens(user_id: str, amount: int, description: str, user_email: str = None) -> bool:
    # Админ не платит токены
    if user_email and is_admin(user_email):
        return True
    
    user = await db.users.find_one({'id': user_id})
    if not user or user.get('balance', 0) < amount:
        return False
    
    await db.users.update_one(
        {'id': user_id},
        {
            '$inc': {'balance': -amount},
            '$push': {
                'transactions': {
                    'id': str(uuid.uuid4()),
                    'amount': -amount,
                    'description': description,
                    'created_at': datetime.now(timezone.utc).isoformat()
                }
            }
        }
    )
    return True

# ==================== FILE PROCESSING ====================

def extract_excel_content(file_path: str) -> str:
    try:
        if file_path.endswith('.xls'):
            workbook = xlrd.open_workbook(file_path)
            content = []
            for sheet in workbook.sheets():
                content.append(f"\n=== Лист: {sheet.name} ===")
                for row_idx in range(min(sheet.nrows, 100)):
                    row = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
                    content.append(" | ".join(row))
            return "\n".join(content)
        else:
            workbook = openpyxl.load_workbook(file_path)
            content = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content.append(f"\n=== Лист: {sheet_name} ===")
                row_count = 0
                for row in sheet.iter_rows(values_only=True):
                    if row_count >= 100:
                        break
                    row_str = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    content.append(row_str)
                    row_count += 1
            return "\n".join(content)
    except Exception as e:
        logger.error(f"Error extracting Excel content: {e}")
        return f"Ошибка чтения Excel файла: {str(e)}"

def extract_word_content(file_path: str) -> str:
    try:
        doc = Document(file_path)
        content = []
        for para in doc.paragraphs:
            if para.text.strip():
                content.append(para.text)
        for table in doc.tables:
            content.append("\n[Таблица]")
            for row in table.rows:
                row_text = " | ".join([cell.text for cell in row.cells])
                content.append(row_text)
        return "\n".join(content)
    except Exception as e:
        logger.error(f"Error extracting Word content: {e}")
        return f"Ошибка чтения Word файла: {str(e)}"

def extract_powerpoint_content(file_path: str) -> str:
    try:
        prs = Presentation(file_path)
        content = []
        content.append(f"Презентация содержит {len(prs.slides)} слайдов.")
        content.append(f"Размер слайдов: {prs.slide_width.pt:.0f}x{prs.slide_height.pt:.0f} pt\n")
        
        for idx, slide in enumerate(prs.slides, 1):
            content.append(f"\n=== Слайд {idx} ===")
            layout_name = slide.slide_layout.name if slide.slide_layout else "Без макета"
            content.append(f"Макет: {layout_name}")
            
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(f"  - {shape.text[:200]}")
            if texts:
                content.append("Текст на слайде:")
                content.extend(texts)
        return "\n".join(content)
    except Exception as e:
        logger.error(f"Error extracting PowerPoint content: {e}")
        return f"Ошибка чтения PowerPoint файла: {str(e)}"

def extract_pdf_content(file_path: str) -> str:
    try:
        content = []
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            content.append(f"PDF документ содержит {len(reader.pages)} страниц.\n")
            for i, page in enumerate(reader.pages[:20], 1):
                text = page.extract_text()
                if text:
                    content.append(f"\n=== Страница {i} ===")
                    content.append(text[:3000])
        return "\n".join(content)
    except Exception as e:
        logger.error(f"Error extracting PDF content: {e}")
        return f"Ошибка чтения PDF файла: {str(e)}"

def extract_txt_content(file_path: str) -> str:
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()[:20000]
    except UnicodeDecodeError:
        with open(file_path, 'r', encoding='latin-1') as file:
            return file.read()[:20000]
    except Exception as e:
        return f"Ошибка чтения текстового файла: {str(e)}"

def extract_file_content(file_path: str, file_type: str) -> str:
    extractors = {
        'xlsx': extract_excel_content,
        'xls': extract_excel_content,
        'docx': extract_word_content,
        'pptx': extract_powerpoint_content,
        'pdf': extract_pdf_content,
        'txt': extract_txt_content,
        'rtf': extract_txt_content,
    }
    extractor = extractors.get(file_type.lower())
    if extractor:
        return extractor(file_path)
    return "Неподдерживаемый тип файла"

def get_file_type(filename: str) -> str:
    return filename.lower().split('.')[-1]

# ==================== DOCUMENT GENERATION ====================

def generate_excel_from_data(data: dict, filename: str) -> str:
    wb = Workbook()
    ws = wb.active
    
    if 'title' in data:
        ws.title = data['title'][:31]
    
    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="4F46E5", end_color="4F46E5", fill_type="solid")
    header_alignment = Alignment(horizontal="center", vertical="center")
    
    if 'headers' in data:
        for col, header in enumerate(data['headers'], 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = header_alignment
    
    if 'rows' in data:
        start_row = 2 if 'headers' in data else 1
        for row_idx, row_data in enumerate(data['rows'], start_row):
            for col_idx, cell_value in enumerate(row_data, 1):
                ws.cell(row=row_idx, column=col_idx, value=cell_value)
    
    for column in ws.columns:
        max_length = 0
        column_letter = column[0].column_letter
        for cell in column:
            try:
                if len(str(cell.value)) > max_length:
                    max_length = len(str(cell.value))
            except:
                pass
        ws.column_dimensions[column_letter].width = min(max_length + 2, 50)
    
    file_path = GENERATED_DIR / filename
    wb.save(file_path)
    return str(file_path)

def generate_word_from_data(data: dict, filename: str) -> str:
    doc = Document()
    
    if 'title' in data:
        title = doc.add_heading(data['title'], 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    if 'sections' in data:
        for section in data['sections']:
            if 'heading' in section:
                doc.add_heading(section['heading'], level=section.get('level', 1))
            if 'content' in section:
                doc.add_paragraph(section['content'])
            if 'bullets' in section:
                for bullet in section['bullets']:
                    doc.add_paragraph(bullet, style='List Bullet')
    
    if 'content' in data and isinstance(data['content'], str):
        doc.add_paragraph(data['content'])
    
    if 'table' in data:
        table_data = data['table']
        if 'headers' in table_data and 'rows' in table_data:
            table = doc.add_table(rows=1, cols=len(table_data['headers']))
            table.style = 'Table Grid'
            header_cells = table.rows[0].cells
            for i, header in enumerate(table_data['headers']):
                header_cells[i].text = str(header)
            for row_data in table_data['rows']:
                row_cells = table.add_row().cells
                for i, cell_value in enumerate(row_data):
                    row_cells[i].text = str(cell_value)
    
    file_path = GENERATED_DIR / filename
    doc.save(file_path)
    return str(file_path)

def generate_pptx_from_template(data: dict, filename: str, template_path: str = None) -> str:
    """Generate PowerPoint using template if available"""
    
    if template_path and os.path.exists(template_path):
        # Используем шаблон
        prs = Presentation(template_path)
        logger.info(f"Using template: {template_path}")
        
        slides_data = data.get('slides', [])
        
        # Заполняем существующие слайды или добавляем новые
        for i, slide_data in enumerate(slides_data):
            if i < len(prs.slides):
                # Заполняем существующий слайд
                slide = prs.slides[i]
                _fill_slide_content(slide, slide_data)
            else:
                # Добавляем новый слайд с макетом из шаблона
                if len(prs.slide_layouts) > 1:
                    layout = prs.slide_layouts[1]  # Content layout
                else:
                    layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(layout)
                _fill_slide_content(slide, slide_data)
    else:
        # Создаём с нуля
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)
        
        slides_data = data.get('slides', [])
        
        for i, slide_data in enumerate(slides_data):
            slide_type = slide_data.get('type', 'content')
            
            if slide_type == 'title' or i == 0:
                layout = prs.slide_layouts[0]  # Title slide
            else:
                layout = prs.slide_layouts[1]  # Content slide
            
            slide = prs.slides.add_slide(layout)
            _fill_slide_content(slide, slide_data)
    
    file_path = GENERATED_DIR / filename
    prs.save(file_path)
    return str(file_path)

def _fill_slide_content(slide, slide_data: dict):
    """Fill slide with content from data"""
    title_text = slide_data.get('title', '')
    subtitle_text = slide_data.get('subtitle', '')
    content = slide_data.get('content', [])
    bullets = slide_data.get('bullets', [])
    
    # Заполняем заголовок
    if slide.shapes.title and title_text:
        slide.shapes.title.text = title_text
    
    # Ищем placeholder для контента
    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape == slide.shapes.title:
                continue
            
            tf = shape.text_frame
            
            # Если это subtitle placeholder
            if subtitle_text and hasattr(shape, 'placeholder_format'):
                if shape.placeholder_format and shape.placeholder_format.type:
                    placeholder_type = shape.placeholder_format.type
                    # Subtitle placeholder type is usually 2
                    if placeholder_type in [2, 4]:
                        tf.clear()
                        tf.paragraphs[0].text = subtitle_text
                        continue
            
            # Заполняем контент
            if bullets:
                tf.clear()
                for j, bullet in enumerate(bullets):
                    if j == 0:
                        tf.paragraphs[0].text = str(bullet)
                    else:
                        p = tf.add_paragraph()
                        p.text = str(bullet)
                        p.level = 0
            elif content:
                tf.clear()
                if isinstance(content, list):
                    for j, item in enumerate(content):
                        if j == 0:
                            tf.paragraphs[0].text = str(item)
                        else:
                            p = tf.add_paragraph()
                            p.text = str(item)
                else:
                    tf.paragraphs[0].text = str(content)

# ==================== AUTH ROUTES ====================

@api_router.post("/auth/register", response_model=TokenResponse)
async def register(user_data: UserCreate):
    existing = await db.users.find_one({'email': user_data.email.lower()})
    if existing:
        raise HTTPException(status_code=400, detail='Email уже зарегистрирован')
    
    if not user_data.agree_terms:
        raise HTTPException(status_code=400, detail='Необходимо согласиться с условиями')
    
    user_id = str(uuid.uuid4())
    user_email = user_data.email.lower()
    user_is_admin = is_admin(user_email)
    
    user = {
        'id': user_id,
        'email': user_email,
        'password_hash': hash_password(user_data.password),
        'balance': 999999 if user_is_admin else INITIAL_TOKENS,
        'is_admin': user_is_admin,
        'is_blocked': False,
        'agreed_terms': True,
        'agreed_at': datetime.now(timezone.utc).isoformat(),
        'transactions': [{
            'id': str(uuid.uuid4()),
            'amount': 999999 if user_is_admin else INITIAL_TOKENS,
            'description': 'Админский аккаунт' if user_is_admin else f'Бонус за регистрацию ({INITIAL_TOKENS} токенов)',
            'created_at': datetime.now(timezone.utc).isoformat()
        }],
        'created_at': datetime.now(timezone.utc).isoformat()
    }
    await db.users.insert_one(user)
    
    token = create_token(user_id)
    
    return TokenResponse(
        access_token=token,
        user=UserResponse(
            id=user_id,
            email=user_email,
            balance=user['balance'],
            is_admin=user_is_admin,
            created_at=datetime.fromisoformat(user['created_at'])
        )
    )

@api_router.post("/auth/login", response_model=TokenResponse)
async def login(user_data: UserLogin):
    user = await db.users.find_one({'email': user_data.email.lower()})
    if not user or not verify_password(user_data.password, user['password_hash']):
        raise HTTPException(status_code=401, detail='Неверный email или пароль')
    
    if user.get('is_blocked'):
        raise HTTPException(status_code=403, detail='Аккаунт заблокирован')
    
    token = create_token(user['id'])
    
    return TokenResponse(
        access_token=token,
        user=UserResponse(
            id=user['id'],
            email=user['email'],
            balance=user.get('balance', 0),
            is_admin=user.get('is_admin', False),
            created_at=datetime.fromisoformat(user['created_at']) if isinstance(user['created_at'], str) else user['created_at']
        )
    )

@api_router.get("/auth/me", response_model=UserResponse)
async def get_me(user = Depends(get_current_user)):
    return UserResponse(
        id=user['id'],
        email=user['email'],
        balance=user.get('balance', 0),
        is_admin=user.get('is_admin', False),
        created_at=datetime.fromisoformat(user['created_at']) if isinstance(user['created_at'], str) else user['created_at']
    )

# ==================== BALANCE ROUTES ====================

@api_router.get("/balance", response_model=BalanceResponse)
async def get_balance(user = Depends(get_current_user)):
    return BalanceResponse(
        balance=user.get('balance', 0),
        token_costs=TOKEN_COSTS
    )

@api_router.post("/balance/topup")
async def topup_balance(request: TopUpRequest, user = Depends(get_current_user)):
    if request.amount < 1:
        raise HTTPException(status_code=400, detail='Минимальная сумма: 1 рубль')
    
    tokens = request.amount
    
    await db.users.update_one(
        {'id': user['id']},
        {
            '$inc': {'balance': tokens},
            '$push': {
                'transactions': {
                    'id': str(uuid.uuid4()),
                    'amount': tokens,
                    'description': f'Пополнение на {request.amount} руб.',
                    'created_at': datetime.now(timezone.utc).isoformat()
                }
            }
        }
    )
    
    updated_user = await db.users.find_one({'id': user['id']})
    return {'success': True, 'new_balance': updated_user.get('balance', 0), 'tokens_added': tokens}

@api_router.get("/balance/history")
async def get_transaction_history(user = Depends(get_current_user)):
    transactions = user.get('transactions', [])
    return {'transactions': transactions[-50:][::-1]}

# ==================== ADMIN ROUTES ====================

@api_router.get("/admin/users")
async def admin_get_users(user = Depends(get_admin_user)):
    users = await db.users.find({}, {'_id': 0, 'password_hash': 0}).sort('created_at', -1).to_list(1000)
    return {'users': users}

@api_router.get("/admin/stats")
async def admin_get_stats(user = Depends(get_admin_user)):
    total_users = await db.users.count_documents({})
    total_chats = await db.chats.count_documents({})
    total_files = await db.files.count_documents({})
    generated_files = await db.files.count_documents({'is_generated': True})
    
    return {
        'total_users': total_users,
        'total_chats': total_chats,
        'total_files': total_files,
        'generated_files': generated_files
    }

@api_router.patch("/admin/users/{user_id}")
async def admin_update_user(user_id: str, update_data: AdminUserUpdate, admin = Depends(get_admin_user)):
    update_fields = {}
    
    if update_data.balance is not None:
        update_fields['balance'] = update_data.balance
    if update_data.is_blocked is not None:
        update_fields['is_blocked'] = update_data.is_blocked
    
    if not update_fields:
        raise HTTPException(status_code=400, detail='Нет данных для обновления')
    
    result = await db.users.update_one({'id': user_id}, {'$set': update_fields})
    if result.modified_count == 0:
        raise HTTPException(status_code=404, detail='Пользователь не найден')
    
    return {'success': True}

@api_router.post("/admin/users/{user_id}/add-tokens")
async def admin_add_tokens(user_id: str, amount: int = Query(...), admin = Depends(get_admin_user)):
    await db.users.update_one(
        {'id': user_id},
        {
            '$inc': {'balance': amount},
            '$push': {
                'transactions': {
                    'id': str(uuid.uuid4()),
                    'amount': amount,
                    'description': f'Начислено администратором',
                    'created_at': datetime.now(timezone.utc).isoformat()
                }
            }
        }
    )
    return {'success': True}

# ==================== CHAT ROUTES ====================

@api_router.post("/chat/create", response_model=ChatResponse)
async def create_chat(chat_data: ChatCreate, user = Depends(get_current_user)):
    chat_id = str(uuid.uuid4())
    now = datetime.now(timezone.utc).isoformat()
    
    chat = {
        'id': chat_id,
        'user_id': user['id'],
        'title': chat_data.title or 'Новый чат',
        'created_at': now,
        'updated_at': now
    }
    await db.chats.insert_one(chat)
    
    return ChatResponse(
        id=chat_id,
        user_id=user['id'],
        title=chat['title'],
        created_at=datetime.fromisoformat(now),
        updated_at=datetime.fromisoformat(now)
    )

@api_router.get("/chat/list", response_model=List[ChatResponse])
async def list_chats(user = Depends(get_current_user)):
    chats = await db.chats.find({'user_id': user['id']}, {'_id': 0}).sort('updated_at', -1).to_list(100)
    return [
        ChatResponse(
            id=c['id'],
            user_id=c['user_id'],
            title=c['title'],
            created_at=datetime.fromisoformat(c['created_at']) if isinstance(c['created_at'], str) else c['created_at'],
            updated_at=datetime.fromisoformat(c['updated_at']) if isinstance(c['updated_at'], str) else c['updated_at']
        )
        for c in chats
    ]

@api_router.get("/chat/{chat_id}", response_model=ChatDetailResponse)
async def get_chat(chat_id: str, user = Depends(get_current_user)):
    chat = await db.chats.find_one({'id': chat_id, 'user_id': user['id']})
    if not chat:
        raise HTTPException(status_code=404, detail='Чат не найден')
    
    messages = await db.messages.find({'chat_id': chat_id}, {'_id': 0}).sort('created_at', 1).to_list(1000)
    files = await db.files.find({'chat_id': chat_id}, {'_id': 0}).sort('created_at', -1).to_list(100)
    
    return ChatDetailResponse(
        id=chat['id'],
        user_id=chat['user_id'],
        title=chat['title'],
        messages=[
            MessageResponse(
                id=m['id'],
                chat_id=m['chat_id'],
                role=m['role'],
                content=m['content'],
                file_id=m.get('file_id'),
                file_name=m.get('file_name'),
                tokens_used=m.get('tokens_used'),
                created_at=datetime.fromisoformat(m['created_at']) if isinstance(m['created_at'], str) else m['created_at']
            )
            for m in messages
        ],
        files=[
            FileResponse_(
                id=f['id'],
                filename=f['filename'],
                file_type=f['file_type'],
                is_generated=f.get('is_generated', False),
                created_at=datetime.fromisoformat(f['created_at']) if isinstance(f['created_at'], str) else f['created_at']
            )
            for f in files
        ],
        created_at=datetime.fromisoformat(chat['created_at']) if isinstance(chat['created_at'], str) else chat['created_at'],
        updated_at=datetime.fromisoformat(chat['updated_at']) if isinstance(chat['updated_at'], str) else chat['updated_at']
    )

@api_router.delete("/chat/{chat_id}")
async def delete_chat(chat_id: str, user = Depends(get_current_user)):
    chat = await db.chats.find_one({'id': chat_id, 'user_id': user['id']})
    if not chat:
        raise HTTPException(status_code=404, detail='Чат не найден')
    
    files = await db.files.find({'chat_id': chat_id}).to_list(100)
    for f in files:
        try:
            if os.path.exists(f['file_path']):
                os.remove(f['file_path'])
        except:
            pass
    
    await db.files.delete_many({'chat_id': chat_id})
    await db.messages.delete_many({'chat_id': chat_id})
    await db.chats.delete_one({'id': chat_id})
    
    return {'status': 'deleted'}

# ==================== MESSAGE ROUTES ====================

@api_router.post("/chat/{chat_id}/upload")
async def upload_file(chat_id: str, file: UploadFile = File(...), user = Depends(get_current_user)):
    chat = await db.chats.find_one({'id': chat_id, 'user_id': user['id']})
    if not chat:
        raise HTTPException(status_code=404, detail='Чат не найден')
    
    file_type = get_file_type(file.filename)
    allowed_types = ['xlsx', 'xls', 'docx', 'pptx', 'pdf', 'txt', 'rtf']
    if file_type not in allowed_types:
        raise HTTPException(status_code=400, detail=f'Тип файла не поддерживается. Разрешены: {allowed_types}')
    
    file_id = str(uuid.uuid4())
    file_path = UPLOAD_DIR / f"{file_id}_{file.filename}"
    
    with open(file_path, 'wb') as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    extracted_content = extract_file_content(str(file_path), file_type)
    
    now = datetime.now(timezone.utc).isoformat()
    file_record = {
        'id': file_id,
        'chat_id': chat_id,
        'user_id': user['id'],
        'filename': file.filename,
        'file_type': file_type,
        'file_path': str(file_path),
        'extracted_content': extracted_content[:50000],
        'created_at': now,
        'is_generated': False
    }
    await db.files.insert_one(file_record)
    
    msg_id = str(uuid.uuid4())
    message = {
        'id': msg_id,
        'chat_id': chat_id,
        'role': 'user',
        'content': f'Загружен файл: {file.filename}',
        'file_id': file_id,
        'file_name': file.filename,
        'created_at': now
    }
    await db.messages.insert_one(message)
    
    messages_count = await db.messages.count_documents({'chat_id': chat_id})
    if messages_count == 1:
        await db.chats.update_one({'id': chat_id}, {'$set': {'title': f'Работа с: {file.filename}', 'updated_at': now}})
    else:
        await db.chats.update_one({'id': chat_id}, {'$set': {'updated_at': now}})
    
    return {
        'file_id': file_id,
        'filename': file.filename,
        'file_type': file_type,
        'message_id': msg_id,
        'extracted_preview': extracted_content[:500] + '...' if len(extracted_content) > 500 else extracted_content
    }

@api_router.post("/chat/{chat_id}/message", response_model=MessageResponse)
async def send_message(chat_id: str, message_data: MessageCreate, user = Depends(get_current_user)):
    chat = await db.chats.find_one({'id': chat_id, 'user_id': user['id']})
    if not chat:
        raise HTTPException(status_code=404, detail='Чат не найден')
    
    user_is_admin = user.get('is_admin', False)
    
    if not user_is_admin and user.get('balance', 0) < TOKEN_COSTS['analyze']:
        raise HTTPException(
            status_code=402, 
            detail=f'Недостаточно токенов. Минимум нужно: {TOKEN_COSTS["analyze"]} токенов. Ваш баланс: {user.get("balance", 0)}'
        )
    
    now = datetime.now(timezone.utc).isoformat()
    
    user_msg_id = str(uuid.uuid4())
    user_message = {
        'id': user_msg_id,
        'chat_id': chat_id,
        'role': 'user',
        'content': message_data.content,
        'created_at': now
    }
    await db.messages.insert_one(user_message)
    
    files = await db.files.find({'chat_id': chat_id}, {'_id': 0}).to_list(10)
    messages = await db.messages.find({'chat_id': chat_id}, {'_id': 0}).sort('created_at', -1).to_list(20)
    messages.reverse()
    
    # Ищем шаблон презентации
    pptx_template_path = None
    file_context = ""
    
    for f in files:
        file_context += f"\n\n=== Файл: {f['filename']} (тип: {f['file_type']}) ===\n{f.get('extracted_content', '')[:15000]}"
        if f['file_type'] == 'pptx' and not f.get('is_generated'):
            pptx_template_path = f['file_path']
    
    system_message = """Ты - ассистент для работы с документами. Ты помогаешь пользователям:
1. Анализировать загруженные файлы (Excel, Word, PowerPoint, PDF, TXT)
2. Создавать новые документы на основе шаблонов и данных
3. ВАЖНО: Если загружен шаблон презентации (.pptx), ты ОБЯЗАН использовать его структуру!

Когда нужно СОЗДАТЬ или ИЗМЕНИТЬ документ, верни JSON:

Для Excel:
```json
{"action": "create_excel", "is_edit": false, "data": {"title": "Название", "headers": ["Колонка1", "Колонка2"], "rows": [["значение1", "значение2"]]}}
```

Для Word:
```json
{"action": "create_word", "is_edit": false, "data": {"title": "Заголовок", "sections": [{"heading": "Раздел 1", "content": "Текст"}, {"heading": "Раздел 2", "bullets": ["Пункт 1", "Пункт 2"]}]}}
```

Для PowerPoint презентации:
```json
{"action": "create_pptx", "is_edit": false, "use_template": true, "data": {"slides": [{"type": "title", "title": "Заголовок презентации", "subtitle": "Подзаголовок"}, {"type": "content", "title": "Слайд 2", "bullets": ["Пункт 1", "Пункт 2", "Пункт 3"]}, {"type": "content", "title": "Слайд 3", "content": ["Текст первого абзаца", "Текст второго абзаца"]}]}}
```

Поле "is_edit": true если это ДОРАБОТКА существующего документа.
Поле "use_template": true если нужно использовать загруженный шаблон презентации.

Если пользователь загрузил шаблон презентации и PDF с данными - используй данные из PDF для заполнения презентации ПО СТРУКТУРЕ ШАБЛОНА.

Контекст загруженных файлов:
""" + file_context
    
    try:
        chat_session = LlmChat(
            api_key=EMERGENT_LLM_KEY,
            session_id=f"chat_{chat_id}_{uuid.uuid4()}",
            system_message=system_message
        ).with_model("openai", "gpt-4o")
        
        conversation = ""
        for msg in messages[-10:]:
            role = "Пользователь" if msg['role'] == 'user' else "Ассистент"
            conversation += f"\n{role}: {msg['content']}"
        
        conversation += f"\nПользователь: {message_data.content}"
        
        ai_response = await chat_session.send_message(UserMessage(text=conversation))
        
    except Exception as e:
        logger.error(f"AI Error: {e}")
        ai_response = f"Извините, произошла ошибка при обработке запроса: {str(e)}"
    
    generated_file_id = None
    generated_file_name = None
    tokens_used = TOKEN_COSTS['analyze']
    
    json_str = None
    if '```json' in ai_response:
        try:
            json_start = ai_response.find('```json') + 7
            json_end = ai_response.find('```', json_start)
            json_str = ai_response[json_start:json_end].strip()
        except:
            pass
    elif ai_response.strip().startswith('{') and '"action"' in ai_response:
        json_str = ai_response.strip()
    
    if json_str:
        try:
            action_data = json.loads(json_str)
            action = action_data.get('action', '')
            is_edit = action_data.get('is_edit', False)
            use_template = action_data.get('use_template', False)
            data = action_data.get('data', {})
            
            generated_file_id = str(uuid.uuid4())
            file_path = None
            
            if action == 'create_excel':
                tokens_used = TOKEN_COSTS['xlsx_edit'] if is_edit else TOKEN_COSTS['xlsx_create']
                generated_filename = f"generated_{generated_file_id}.xlsx"
                file_path = generate_excel_from_data(data, generated_filename)
                generated_file_name = data.get('title', 'Таблица') + '.xlsx'
                
            elif action == 'create_word':
                tokens_used = TOKEN_COSTS['docx_edit'] if is_edit else TOKEN_COSTS['docx_create']
                generated_filename = f"generated_{generated_file_id}.docx"
                file_path = generate_word_from_data(data, generated_filename)
                generated_file_name = data.get('title', 'Документ') + '.docx'
                
            elif action == 'create_pptx':
                tokens_used = TOKEN_COSTS['pptx_edit'] if is_edit else TOKEN_COSTS['pptx_create']
                generated_filename = f"generated_{generated_file_id}.pptx"
                
                # Используем шаблон если он есть и запрошено use_template
                template = pptx_template_path if use_template else None
                file_path = generate_pptx_from_template(data, generated_filename, template)
                
                # Название файла
                if data.get('slides'):
                    first_slide = data['slides'][0]
                    generated_file_name = first_slide.get('title', 'Презентация')[:50] + '.pptx'
                else:
                    generated_file_name = 'Презентация.pptx'
            
            if file_path and generated_file_name:
                if not user_is_admin and user.get('balance', 0) < tokens_used:
                    raise HTTPException(
                        status_code=402,
                        detail=f'Недостаточно токенов. Нужно: {tokens_used}, ваш баланс: {user.get("balance", 0)}'
                    )
                
                file_record = {
                    'id': generated_file_id,
                    'chat_id': chat_id,
                    'user_id': user['id'],
                    'filename': generated_file_name,
                    'file_type': generated_filename.split('.')[-1],
                    'file_path': file_path,
                    'extracted_content': '',
                    'created_at': now,
                    'is_generated': True
                }
                await db.files.insert_one(file_record)
                
                action_type = 'доработка' if is_edit else 'создание'
                file_type_name = {'xlsx': 'Excel', 'docx': 'Word', 'pptx': 'PowerPoint'}.get(generated_filename.split('.')[-1], 'файла')
                template_note = ' (на основе вашего шаблона)' if use_template and pptx_template_path else ''
                
                if '```json' in ai_response:
                    ai_response = ai_response[:ai_response.find('```json')].strip()
                else:
                    ai_response = ""
                
                token_note = '' if user_is_admin else f'\n💰 Списано {tokens_used} токенов ({action_type} {file_type_name}).'
                ai_response += f"\n\n✅ Файл '{generated_file_name}' создан{template_note}! Вы можете скачать его ниже.{token_note}"
                
        except json.JSONDecodeError as e:
            logger.error(f"JSON parse error: {e}")
        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"Error creating file: {e}")
            generated_file_id = None
            generated_file_name = None
    
    if not user_is_admin:
        await deduct_tokens(user['id'], tokens_used, f"Чат: {chat['title'][:30]}", user.get('email'))
    
    ai_msg_id = str(uuid.uuid4())
    ai_message = {
        'id': ai_msg_id,
        'chat_id': chat_id,
        'role': 'assistant',
        'content': ai_response,
        'file_id': generated_file_id,
        'file_name': generated_file_name,
        'tokens_used': tokens_used if not user_is_admin else 0,
        'created_at': datetime.now(timezone.utc).isoformat()
    }
    await db.messages.insert_one(ai_message)
    
    await db.chats.update_one({'id': chat_id}, {'$set': {'updated_at': datetime.now(timezone.utc).isoformat()}})
    
    return MessageResponse(
        id=ai_msg_id,
        chat_id=chat_id,
        role='assistant',
        content=ai_response,
        file_id=generated_file_id,
        file_name=generated_file_name,
        tokens_used=tokens_used if not user_is_admin else 0,
        created_at=datetime.fromisoformat(ai_message['created_at'])
    )

# ==================== FILE DOWNLOAD ====================

@api_router.get("/files/{file_id}/download")
async def download_file(file_id: str, token: Optional[str] = None):
    if not token:
        raise HTTPException(status_code=401, detail='Требуется авторизация')
    
    try:
        payload = jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
        user_id = payload.get('user_id')
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail='Токен истёк')
    except:
        raise HTTPException(status_code=401, detail='Неверный токен')
    
    file_record = await db.files.find_one({'id': file_id})
    if not file_record:
        raise HTTPException(status_code=404, detail='Файл не найден')
    
    chat = await db.chats.find_one({'id': file_record['chat_id'], 'user_id': user_id})
    if not chat:
        raise HTTPException(status_code=403, detail='Нет доступа к файлу')
    
    if not os.path.exists(file_record['file_path']):
        raise HTTPException(status_code=404, detail='Файл не найден на диске')
    
    mime_types = {
        'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'xls': 'application/vnd.ms-excel',
        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'pdf': 'application/pdf',
        'txt': 'text/plain',
    }
    file_ext = file_record['filename'].split('.')[-1].lower()
    media_type = mime_types.get(file_ext, 'application/octet-stream')
    
    encoded_filename = quote(file_record['filename'])
    
    return FileResponse(
        path=file_record['file_path'],
        filename=file_record['filename'],
        media_type=media_type,
        headers={'Content-Disposition': f"attachment; filename*=UTF-8''{encoded_filename}"}
    )

# ==================== LEGAL PAGES ====================

@api_router.get("/legal/privacy")
async def get_privacy_policy():
    return {
        'title': 'Политика конфиденциальности',
        'content': '''Политика конфиденциальности сервиса DocAI Chat

1. ОБЩИЕ ПОЛОЖЕНИЯ
1.1. Настоящая Политика конфиденциальности определяет порядок обработки персональных данных пользователей сервиса DocAI Chat.
1.2. Используя сервис, вы соглашаетесь с условиями данной Политики.

2. СОБИРАЕМЫЕ ДАННЫЕ
2.1. Мы собираем следующие данные:
- Адрес электронной почты
- Загружаемые документы (обрабатываются для предоставления услуг)
- История использования сервиса

3. ИСПОЛЬЗОВАНИЕ ДАННЫХ
3.1. Ваши данные используются для:
- Предоставления услуг сервиса
- Обработки документов с помощью ИИ
- Связи с вами по вопросам сервиса

4. ЗАЩИТА ДАННЫХ
4.1. Мы применяем современные методы защиты данных.
4.2. Загруженные документы хранятся в зашифрованном виде.
4.3. Доступ к данным имеют только авторизованные сотрудники.

5. ПЕРЕДАЧА ДАННЫХ ТРЕТЬИМ ЛИЦАМ
5.1. Мы не передаём ваши персональные данные третьим лицам без вашего согласия.
5.2. Исключение: требования законодательства РФ.

6. УДАЛЕНИЕ ДАННЫХ
6.1. Вы можете запросить удаление своих данных, написав на email администратора.

7. КОНТАКТЫ
По вопросам конфиденциальности: mishkasladkolapka@gmail.com'''
    }

@api_router.get("/legal/terms")
async def get_terms_of_service():
    return {
        'title': 'Пользовательское соглашение',
        'content': '''Пользовательское соглашение сервиса DocAI Chat

1. ОБЩИЕ ПОЛОЖЕНИЯ
1.1. Настоящее Соглашение регулирует отношения между пользователем и сервисом DocAI Chat.
1.2. Регистрируясь в сервисе, вы принимаете условия данного Соглашения.

2. ОПИСАНИЕ СЕРВИСА
2.1. DocAI Chat — сервис для работы с документами с использованием искусственного интеллекта.
2.2. Сервис позволяет создавать и редактировать презентации, таблицы и текстовые документы.

3. СТОИМОСТЬ УСЛУГ
3.1. Оплата производится в токенах (1 рубль = 1 токен).
3.2. При регистрации предоставляется 300 бесплатных токенов.
3.3. Тарифы:
- Создание презентации: 65 токенов
- Доработка презентации: 10 токенов
- Создание таблицы Excel: 40 токенов
- Доработка таблицы: 7 токенов
- Создание документа Word: 35 токенов
- Доработка документа: 6 токенов
- Анализ файла / вопрос: 5 токенов

4. ПРАВА И ОБЯЗАННОСТИ
4.1. Пользователь обязуется не использовать сервис для незаконных целей.
4.2. Пользователь несёт ответственность за содержимое загружаемых файлов.
4.3. Администрация вправе заблокировать аккаунт при нарушении правил.

5. ОГРАНИЧЕНИЕ ОТВЕТСТВЕННОСТИ
5.1. Сервис предоставляется "как есть".
5.2. Администрация не несёт ответственности за результаты работы ИИ.
5.3. Рекомендуем проверять сгенерированные документы перед использованием.

6. ВОЗВРАТ СРЕДСТВ
6.1. Возврат неиспользованных токенов возможен в течение 14 дней с момента покупки.
6.2. Для возврата свяжитесь с администрацией.

7. ИЗМЕНЕНИЕ УСЛОВИЙ
7.1. Администрация вправе изменять условия Соглашения.
7.2. Актуальная версия всегда доступна в сервисе.

8. КОНТАКТЫ
Email: mishkasladkolapka@gmail.com'''
    }

# ==================== PRICING ====================

@api_router.get("/pricing")
async def get_pricing():
    return {
        'info': '1 рубль = 1 токен',
        'initial_tokens': INITIAL_TOKENS,
        'costs': {
            'Создание презентации PowerPoint': TOKEN_COSTS['pptx_create'],
            'Доработка презентации': TOKEN_COSTS['pptx_edit'],
            'Создание документа Word': TOKEN_COSTS['docx_create'],
            'Доработка документа Word': TOKEN_COSTS['docx_edit'],
            'Создание таблицы Excel': TOKEN_COSTS['xlsx_create'],
            'Доработка таблицы Excel': TOKEN_COSTS['xlsx_edit'],
            'Анализ файла / вопрос': TOKEN_COSTS['analyze'],
        },
        'example': {
            '400 рублей': 'примерно 6 новых презентаций или 10 таблиц Excel или 11 документов Word'
        }
    }

# ==================== HEALTH ====================

@api_router.get("/")
async def root():
    return {"message": "DocAI Chat API", "status": "running"}

@api_router.get("/health")
async def health():
    return {"status": "healthy"}

app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
